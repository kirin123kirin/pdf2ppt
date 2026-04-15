"""
# pip install numpy pymupdf pillow janome surya-ocr "transformers==4.57.3"
pdf2ppt.py - PDF / クリップボード画像 → PowerPoint 変換ツール (surya 0.17.x対応)
=================================================================================
使い方:
  python pdf2ppt.py <pdf_path> [options]
  python pdf2ppt.py              # 引数なし → クリップボードから PDF パスまたは画像を読み取り

オプション:
  --csv              メタデータCSVを出力する (デフォルト: 出力しない)
  -v, --verbose      進捗を詳細表示する (デフォルト: 最小表示)
  --out-dir <dir>    出力先ディレクトリ (デフォルト: PDFと同じ場所 / 画像時はカレントディレクトリ)

処理フロー:
  1. ページを高解像度でラスタライズ (PyMuPDF) ※PDF の場合のみ
  2. surya RecognitionPredictor でテキスト行を検出・認識
  3. フォントサイズを行の高さから推定、フォント色をピクセルから取得
  4. テキスト領域を周辺背景色で塗りつぶし
  5. 残った非白色領域をクラスタリングしてクロップ
  6. PowerPoint に位置情報を使って配置
"""

import argparse
import base64
import csv
import gc
import fitz
import io
import os
import subprocess
import sys
import tempfile
import numpy as np
from scipy.ndimage import binary_fill_holes
from datetime import datetime
from pathlib import Path
from typing import Any
from PIL import Image
from janome.tokenizer import Tokenizer

# ─────────────────────────────────────────────
# 設定
# ─────────────────────────────────────────────

OCR_LANGS = ["ja", "en"]
SURYA_MAX_WIDTH = 2048        # surya に渡す最大幅(px)
RASTER_TARGET_WIDTH = SURYA_MAX_WIDTH  # ラスタライズ目標幅 = surya上限に合わせる（フォントサイズ計算の整合性のため）
SCREENSHOT_DPI = 96.0         # クリップボード画像の想定解像度 (標準スクリーンDPI)
TEXT_PAD_RATIO = 0.15         # テキスト bbox パディング（行高に対する割合）
BG_MARGIN = 8                 # 背景サンプリングマージン(px)
AUTO_FILL_COLOR = True        # True=周辺背景色で塗りつぶし / False=白固定
FIXED_FILL_COLOR = (255, 255, 255)
TEXT_COLOR_TOL = 0            # テキスト色差閾値（0=全ピクセル置換）
BLOCK_GAP = 10                # 画像クラスタリングギャップ(px)
MIN_BLOCK_AREA = 50           # クロップ対象の最小面積(px²)
NON_WHITE_THR = 238           # 非白判定閾値(0-255)
QUALITY_THR = 0.6             # テキスト品質スコア閾値(0.0-1.0)

_VALID_POS = frozenset({"名詞", "動詞", "形容詞", "副詞", "助詞"})
_tokenizer = Tokenizer()

# ─────────────────────────────────────────────
# クリップボード読み取り
# ─────────────────────────────────────────────


def _read_clipboard() -> str:
    """Windows クリップボードからファイルパスを取得する。
    ファイルエクスプローラーからコピーしたファイル (FileDropList) と
    テキストとしてコピーしたパスの両方に対応する。
    """
    ps_script = (
        "[Console]::OutputEncoding = [System.Text.Encoding]::UTF8; "
        "$files = Get-Clipboard -Format FileDropList; "
        "if ($files) { $files[0].FullName } else { Get-Clipboard }"
    )
    try:
        result = subprocess.run(
            ["powershell", "-NoProfile", "-Command", ps_script],
            capture_output=True, timeout=5
        )
        return result.stdout.decode("utf-8").strip()
    except Exception as e:
        print(f"[ERROR] クリップボード読み取り失敗: {e}")
        return ""


def _read_clipboard_image() -> "Image.Image | None":
    """Windows クリップボードからビットマップ画像を取得する。
    スクリーンショット (PrintScreen, Snipping Tool 等) でコピーした画像に対応。
    画像がなければ None を返す。
    """
    ps_script = (
        "Add-Type -AssemblyName System.Windows.Forms; "
        "Add-Type -AssemblyName System.Drawing; "
        "$img = [System.Windows.Forms.Clipboard]::GetImage(); "
        "if ($img -ne $null) { "
        "    $ms = New-Object System.IO.MemoryStream; "
        "    $img.Save($ms, [System.Drawing.Imaging.ImageFormat]::Png); "
        "    [Convert]::ToBase64String($ms.ToArray()) "
        "}"
    )
    try:
        result = subprocess.run(
            ["powershell", "-NoProfile", "-Command", ps_script],
            capture_output=True, timeout=30
        )
        data = result.stdout.strip()
        if not data:
            return None
        img_bytes = base64.b64decode(data)
        return Image.open(io.BytesIO(img_bytes))
    except Exception as e:
        print(f"[WARN] クリップボード画像読み取り失敗: {e}")
        return None


# ─────────────────────────────────────────────
# ラスタライズ
# ─────────────────────────────────────────────

def _page_dpi(page) -> tuple[float, float]:
    """ページ幅(pt)から目標幅(px)になるDPIを返す。戻り値: (dpi, page_w_pt)"""
    w_pt = page.rect.width
    return RASTER_TARGET_WIDTH / w_pt * 72, w_pt


def rasterize(page) -> tuple[Image.Image, float]:
    """PDFページを自動DPIでラスタライズ。戻り値: (img, dpi)"""
    dpi, _ = _page_dpi(page)
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    pix = page.get_pixmap(matrix=mat, alpha=False)
    # PNG エンコード/デコードを省略して直接 PIL Image に変換
    img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
    return img, dpi


def _resize_for_ocr(img: Image.Image) -> tuple[Image.Image, float]:
    """SURYA_MAX_WIDTH を超える場合にリサイズ。戻り値: (img, scale)"""
    w, h = img.size
    if w <= SURYA_MAX_WIDTH:
        return img, 1.0
    scale = w / SURYA_MAX_WIDTH
    return img.resize((SURYA_MAX_WIDTH, int(h / scale)), Image.LANCZOS), scale


# ─────────────────────────────────────────────
# surya モデルロード
# ─────────────────────────────────────────────

def load_models(verbose: bool = False):
    """surya DetectionPredictor / RecognitionPredictor をロード。"""
    print("surya モデルをロード中...")
    try:
        from surya.recognition import RecognitionPredictor
        from surya.detection import DetectionPredictor
        from surya.foundation import FoundationPredictor
    except ImportError as e:
        print(f"[ERROR] surya のインポートに失敗: {e}\n  pip install surya-ocr")
        sys.exit(1)
    base = FoundationPredictor()
    det = DetectionPredictor()
    rec = RecognitionPredictor(base)
    if verbose:
        print("surya モデルロード完了")
    return det, rec


# ─────────────────────────────────────────────
# テキスト品質スコア
# ─────────────────────────────────────────────

def _ling_score(text: str) -> float:
    """形態素解析による言語品質スコア (0.0-1.0)。"""
    tokens = list(_tokenizer.tokenize(text))
    if not tokens:
        return 0.0
    valid = sum(
        1 for t in tokens
        if t.part_of_speech.split(",")[0] in _VALID_POS
        and getattr(t.node, "node_type", "") != "UNKNOWN"
    )
    noise = sum(
        1 for t in tokens
        if getattr(t.node, "node_type", "") == "UNKNOWN"
        and (not t.extra or t.extra[0] == "*")
    )
    return max(0.0, (valid - noise) / len(tokens))


def score_text(row: dict[str, Any]) -> float:
    """OCR confidence と言語スコアを合成した品質スコア (0.0-1.0)。
    conf >= 0.9 の行は最低スコアが 0.72 > QUALITY_THR なので janome を省略する。
    """
    conf = float(row.get("confidence") or 0.0)
    if conf >= 0.9:
        return conf * 0.8
    ling = _ling_score(row.get("text", ""))
    return conf * 0.6 + ling * 0.4


# ─────────────────────────────────────────────
# STEP 1: テキスト抽出
# ─────────────────────────────────────────────

def _text_color(arr: np.ndarray, x0, y0, x1, y1) -> tuple[int, int, int]:
    """領域内の背景輝度に応じてテキスト色を推定する。
    明背景 → 暗ピクセルの中央値、暗背景 → 明ピクセルの中央値。
    """
    region = arr[y0:y1, x0:x1]
    if region.size == 0:
        return (0, 0, 0)
    lum = np.mean(region, axis=2)
    if lum.mean() < 128:
        mask = lum > 160
        if not mask.any():
            return (255, 255, 255)
    else:
        mask = lum < 160
        if not mask.any():
            return (0, 0, 0)
    med = np.median(region[mask], axis=0).astype(int)
    return tuple(int(v) for v in med)


def _hex(rgb: tuple[int, int, int]) -> str:
    return "#{:02X}{:02X}{:02X}".format(*rgb)


def extract_text(
    page_img: Image.Image,
    dpi: float,
    page_num: int,
    det,
    rec,
    verbose: bool = False,
) -> list[dict[str, Any]]:
    """ページ画像からテキスト行を抽出し、品質フィルタ済みのリストを返す。"""
    ocr_img, scale = _resize_for_ocr(page_img)
    result = rec([ocr_img], det_predictor=det)[0]
    del ocr_img

    arr = np.array(page_img.convert("RGB"))
    h, w = arr.shape[:2]
    rows = []

    for line in result.text_lines:
        text = (line.text or "").strip()
        if not text:
            continue
        bbox = line.bbox
        x0 = max(0, int(bbox[0] * scale))
        y0 = max(0, int(bbox[1] * scale))
        x1 = min(w, int(bbox[2] * scale))
        y1 = min(h, int(bbox[3] * scale))
        if x1 <= x0 or y1 <= y0:
            continue

        conf = getattr(line, "confidence", None)
        row = {
            "page":       page_num + 1,
            "text":       text,
            "font_name":  "unknown",
            "font_size":  _snap_font_size((y1 - y0) / dpi * 72),
            "font_color": _hex(_text_color(arr, x0, y0, x1, y1)),
            "confidence": round(float(conf), 3) if conf is not None else None,
            "x0": round(x0/dpi*72, 2), "y0": round(y0/dpi*72, 2),
            "x1": round(x1/dpi*72, 2), "y1": round(y1/dpi*72, 2),
            "x0_px": x0, "y0_px": y0, "x1_px": x1, "y1_px": y1,
        }
        row["quality"] = round(score_text(row), 3)
        if row["quality"] >= QUALITY_THR:
            rows.append(row)

    return rows


# ─────────────────────────────────────────────
# STEP 2: テキストマスキング
# ─────────────────────────────────────────────

def _sample_bg(arr: np.ndarray, x0, y0, x1, y1) -> tuple[int, int, int]:
    """bbox 外周から背景色を中央値でサンプリング。numpy 配列のまま処理する。"""
    h, w = arr.shape[:2]
    parts = [
        arr[max(0, y0 - BG_MARGIN),  max(0, x0):min(w, x1)],          # 上辺
        arr[min(h - 1, y1 + BG_MARGIN), max(0, x0):min(w, x1)],        # 下辺
        arr[max(0, y0):min(h, y1), max(0, x0 - BG_MARGIN)].reshape(-1, 3),   # 左辺
        arr[max(0, y0):min(h, y1), min(w - 1, x1 + BG_MARGIN)].reshape(-1, 3),  # 右辺
    ]
    parts = [p for p in parts if p.size > 0]
    if not parts:
        return (255, 255, 255)
    return tuple(int(v) for v in np.median(np.vstack(parts), axis=0).astype(int))


def mask_text(page_img: Image.Image, rows: list[dict[str, Any]]) -> Image.Image:
    """テキスト行の bbox を周辺背景色で塗りつぶし、マスク済み画像を返す。
    uint8 のまま処理してメモリ使用量を抑える。
    """
    # float32 ではなく uint8 で保持（メモリ約 1/4）
    base = np.array(page_img.convert("RGB"))   # uint8
    out = base.copy()
    h, w = base.shape[:2]

    if not rows:
        return Image.fromarray(out)

    def pad(row):
        return max(1, int((row["y1_px"] - row["y0_px"]) * TEXT_PAD_RATIO))

    ocr_x0 = max(0, min(row["x0_px"] - pad(row) for row in rows))
    ocr_y0 = max(0, min(row["y0_px"] - pad(row) for row in rows))
    ocr_x1 = min(w, max(row["x1_px"] + pad(row) for row in rows))
    ocr_y1 = min(h, max(row["y1_px"] + pad(row) for row in rows))

    protect = np.ones((h, w), dtype=bool)
    protect[ocr_y0:ocr_y1, ocr_x0:ocr_x1] = False

    for row in rows:
        p = pad(row)
        x0 = max(0, row["x0_px"] - p)
        y0 = max(0, row["y0_px"] - p)
        x1 = min(w, row["x1_px"] + p)
        y1 = min(h, row["y1_px"] + p)

        bg = _sample_bg(base, x0, y0, x1, y1) if AUTO_FILL_COLOR else FIXED_FILL_COLOR
        bg_u8 = np.array(bg, dtype=np.uint8)

        region = out[y0:y1, x0:x1]
        prot = protect[y0:y1, x0:x1]

        if TEXT_COLOR_TOL == 0:
            # 全ピクセル置換（デフォルト）: norm 計算不要
            fill_mask = ~prot
        else:
            diff = region.astype(np.float32) - bg_u8.astype(np.float32)
            fill_mask = (np.linalg.norm(diff, axis=2) > TEXT_COLOR_TOL) & ~prot

        region[fill_mask] = bg_u8

    return Image.fromarray(out)


# ─────────────────────────────────────────────
# STEP 3: 画像ブロック抽出
# ─────────────────────────────────────────────

def _segments(binary_1d: np.ndarray, gap: int) -> list[tuple[int, int]]:
    idx = np.where(binary_1d)[0]
    if len(idx) == 0:
        return []
    segs, start, prev = [], idx[0], idx[0]
    for v in idx[1:]:
        if v - prev > gap:
            segs.append((int(start), int(prev)))
            start = v
        prev = v
    segs.append((int(start), int(prev)))
    return segs


def _merge(blocks: list[tuple], gap: int) -> list[tuple]:
    if not blocks:
        return []
    blocks = sorted(blocks, key=lambda b: (b[1], b[0]))
    merged = [list(blocks[0])]
    for b in blocks[1:]:
        last = merged[-1]
        if (b[0] <= last[2] + gap and b[2] >= last[0] - gap
                and b[1] <= last[3] + gap and b[3] >= last[1] - gap):
            merged[-1] = [min(last[0], b[0]), min(last[1], b[1]),
                          max(last[2], b[2]), max(last[3], b[3])]
        else:
            merged.append(list(b))
    return [tuple(b) for b in merged]


def _detect_blocks(img: Image.Image) -> list[tuple]:
    """非白色領域をクラスタリングして bbox リストを返す。"""
    arr = np.array(img)
    non_white = ~(
        (arr[:, :, 0] > NON_WHITE_THR) &
        (arr[:, :, 1] > NON_WHITE_THR) &
        (arr[:, :, 2] > NON_WHITE_THR)
    )
    if not non_white.any():
        return []
    blocks = []
    for ry0, ry1 in _segments(np.any(non_white, axis=1), BLOCK_GAP):
        for cx0, cx1 in _segments(np.any(non_white[ry0:ry1+1, :], axis=0), BLOCK_GAP):
            if (ry1 - ry0) * (cx1 - cx0) >= MIN_BLOCK_AREA:
                blocks.append((cx0, ry0, cx1, ry1))
    return _merge(blocks, BLOCK_GAP)


def _trim_margins(img: Image.Image) -> tuple[Image.Image, int, int]:
    """非白色領域のタイトな bbox にトリミング。戻り値: (trimmed, dx, dy)"""
    arr = np.array(img)
    nw = ~(
        (arr[:, :, 0] > NON_WHITE_THR) &
        (arr[:, :, 1] > NON_WHITE_THR) &
        (arr[:, :, 2] > NON_WHITE_THR)
    )
    if not nw.any():
        return img, 0, 0
    rows = np.where(np.any(nw, axis=1))[0]
    cols = np.where(np.any(nw, axis=0))[0]
    ty0, ty1 = int(rows[0]), int(rows[-1])
    tx0, tx1 = int(cols[0]), int(cols[-1])
    return img.crop((tx0, ty0, tx1 + 1, ty1 + 1)), tx0, ty0


def _transparent_bg(img: Image.Image) -> Image.Image:
    """外周連結の背景色をアルファ透過にした RGBA 画像を返す。
    BFS の代わりに scipy.ndimage.binary_fill_holes を使い高速化する。
    外周に接する白領域 = is_bg & ~fill_holes(~is_bg)
    """
    arr = np.array(img)   # すでに RGB
    is_bg = (
        (arr[:, :, 0] > NON_WHITE_THR) &
        (arr[:, :, 1] > NON_WHITE_THR) &
        (arr[:, :, 2] > NON_WHITE_THR)
    )
    # ~is_bg の穴を埋めた領域 = 外周非連結の背景も含む → 差分が外周連結背景
    outer = is_bg & ~binary_fill_holes(~is_bg)

    rgba = np.empty((arr.shape[0], arr.shape[1], 4), dtype=np.uint8)
    rgba[:, :, :3] = arr
    rgba[:, :, 3] = 255
    rgba[outer, 3] = 0
    return Image.fromarray(rgba, "RGBA")


def extract_blocks(
    masked_img: Image.Image,
    dpi: float,
    page_num: int,
    verbose: bool = False,
) -> list[dict[str, Any]]:
    """マスク済み画像から図形ブロックを検出・トリミング・透過処理してインメモリで返す。"""
    results = []
    for i, (x0, y0, x1, y1) in enumerate(_detect_blocks(masked_img)):
        trimmed, dx, dy = _trim_margins(masked_img.crop((x0, y0, x1, y1)))
        tx0, ty0 = x0 + dx, y0 + dy
        tx1, ty1 = tx0 + trimmed.width, ty0 + trimmed.height

        buf = io.BytesIO()
        _transparent_bg(trimmed).save(buf, "PNG")
        buf.seek(0)

        results.append({
            "img_bytes": buf,
            "page":      page_num + 1,
            "x0_px": tx0, "y0_px": ty0, "x1_px": tx1, "y1_px": ty1,
            "width_px":  trimmed.width,  "height_px": trimmed.height,
            "x0": round(tx0/dpi*72, 2),  "y0": round(ty0/dpi*72, 2),
            "x1": round(tx1/dpi*72, 2),  "y1": round(ty1/dpi*72, 2),
        })
        if verbose:
            print(f"  画像ブロック検出: {i+1}  ({trimmed.width}x{trimmed.height}px)")
    return results


# ─────────────────────────────────────────────
# STEP 4: PowerPoint 出力（ストリーミング）
# ─────────────────────────────────────────────

_FONT_SIZES = [6, 7, 8, 9, 10, 11, 12, 14, 16, 18, 20,
               24, 28, 32, 36, 40, 44, 48, 54, 60, 66, 72, 80, 88, 96]


def _snap_font_size(pt: float) -> float:
    return min(_FONT_SIZES, key=lambda s: abs(s - pt))


def _init_pptx(page_w_pt: float, page_h_pt: float):
    """Presentation を初期化して返す。"""
    from pptx import Presentation
    from pptx.util import Emu
    prs = Presentation()
    prs.slide_width = Emu(int(page_w_pt * 12700))
    prs.slide_height = Emu(int(page_h_pt * 12700))
    return prs


def _add_pptx_slide(
    prs,
    text_rows: list[dict[str, Any]],
    img_meta: list[dict[str, Any]],
) -> None:
    """1ページ分のスライドを追加する。img_bytes は追加後に即クローズする。"""
    from pptx.util import Emu, Pt
    from pptx.dml.color import RGBColor

    def emu(pt: float) -> Emu:
        return Emu(int(pt * 12700))

    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    for m in img_meta:
        buf = m["img_bytes"]
        buf.seek(0)
        w = max(m["x1"] - m["x0"], 1.0)
        h = max(m["y1"] - m["y0"], 1.0)
        slide.shapes.add_picture(buf, emu(m["x0"]), emu(m["y0"]), emu(w), emu(h))
        buf.close()  # add_picture がデータを読み終えたので即解放

    for r in text_rows:
        w = max(r["x1"] - r["x0"], 1.0)
        h = max(r["y1"] - r["y0"], 1.0)
        tf = slide.shapes.add_textbox(
            emu(r["x0"]), emu(r["y0"]), emu(w), emu(h)
        ).text_frame
        tf.word_wrap = False
        run = tf.paragraphs[0].add_run()
        run.text = r["text"]
        run.font.name = "メイリオ"
        run.font.size = Pt(_snap_font_size(r["font_size"]))
        hc = r.get("font_color", "#000000").lstrip("#")
        if len(hc) == 6:
            run.font.color.rgb = RGBColor(
                int(hc[0:2], 16), int(hc[2:4], 16), int(hc[4:6], 16)
            )


# ─────────────────────────────────────────────
# メタデータ CSV 出力
# ─────────────────────────────────────────────

_META_FIELDS = [
    "type", "page",
    "text", "font_name", "font_size", "font_color", "confidence", "quality",
    "x0", "y0", "x1", "y1", "x0_px", "y0_px", "x1_px", "y1_px",
    "width_px", "height_px",
]


def write_meta_csv(
    text_rows: list[dict[str, Any]],
    img_meta:  list[dict[str, Any]],
    csv_path:  Path,
) -> None:
    """テキスト行と画像ブロックのメタデータを1つのCSVに統合して出力する。"""
    rows = []
    for r in text_rows:
        row = {k: r.get(k, "") for k in _META_FIELDS}
        row["type"] = "text"
        rows.append(row)
    for m in img_meta:
        row = {k: m.get(k, "") for k in _META_FIELDS}
        row["type"] = "image"
        rows.append(row)

    rows.sort(key=lambda r: (r["page"], r.get("y0", 0)))

    with open(csv_path, "w", newline="", encoding="utf-8-sig") as f:
        w = csv.DictWriter(f, fieldnames=_META_FIELDS, extrasaction="ignore")
        w.writeheader()
        w.writerows(rows)
    print(f"メタデータCSV出力完了: {csv_path}  ({len(rows)}件)")


# ─────────────────────────────────────────────
# 処理オーケストレーション
# ─────────────────────────────────────────────

def process_pdf(
    pdf_path: Path,
    out_dir:  Path,
    *,
    pptx_path: Path,
    csv_path:  Path | None,
    verbose:   bool = False,
) -> None:
    """PDF を処理し、PPTX を出力する。csv_path が指定された場合はメタCSVも出力する。"""
    if not pdf_path.exists():
        print(f"[ERROR] ファイルが見つかりません: {pdf_path}")
        sys.exit(1)

    out_dir.mkdir(parents=True, exist_ok=True)

    print(f"処理開始: {pdf_path}")
    if verbose:
        print(f"出力先:   {out_dir}")
        print(f"OCR言語:  {OCR_LANGS}")
        print(f"ラスタライズ目標幅: {RASTER_TARGET_WIDTH}px  surya最大幅: {SURYA_MAX_WIDTH}px")

    det, rec = load_models(verbose=verbose)
    doc = fitz.open(str(pdf_path))
    total = len(doc)

    # 最初のページのサイズで Presentation を初期化
    first_page = doc[0]
    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        print("[ERROR] python-pptx が必要です: pip install python-pptx")
        sys.exit(1)
    prs = _init_pptx(first_page.rect.width, first_page.rect.height)

    csv_texts = [] if csv_path is not None else None
    csv_imgs = [] if csv_path is not None else None

    for page_num, page in enumerate(doc):
        print(f"\nページ {page_num+1}/{total} 処理中...")

        page_img, dpi = rasterize(page)
        if verbose:
            print(f"  ラスタライズ: {page_img.width}x{page_img.height}px  DPI={dpi:.1f}")

        if verbose:
            print("  OCR実行中...")
        rows = extract_text(page_img, dpi, page_num, det, rec, verbose=verbose)
        print(f"  テキスト: {len(rows)}行")

        masked = mask_text(page_img, rows)
        del page_img

        imgs = extract_blocks(masked, dpi, page_num, verbose=verbose)
        print(f"  画像ブロック: {len(imgs)}件")
        del masked

        # スライド追加（BytesIO はここで消費・クローズされる）
        _add_pptx_slide(prs, rows, imgs)

        # CSV指定時のみメタデータを蓄積
        if csv_texts is not None:
            csv_texts.extend(rows)
            csv_imgs.extend({k: v for k, v in m.items() if k != "img_bytes"} for m in imgs)

        del rows, imgs
        gc.collect()

    doc.close()

    prs.save(str(pptx_path))
    print(f"\nPowerPoint出力完了: {pptx_path}  ({total}スライド)")

    if csv_path is not None:
        write_meta_csv(csv_texts, csv_imgs, csv_path)

    print("完了")


def process_image(
    clip_img: "Image.Image",
    out_dir:  Path,
    *,
    pptx_path: Path,
    csv_path:  "Path | None",
    verbose:   bool = False,
) -> None:
    """クリップボード画像を処理し、1スライドの PPTX を出力する。
    画像が RASTER_TARGET_WIDTH を超える場合はリサイズしてメモリを節約する。
    """
    out_dir.mkdir(parents=True, exist_ok=True)

    w_orig, h_orig = clip_img.size

    # RASTER_TARGET_WIDTH を超える場合はリサイズ (surya と同じ上限)
    if w_orig > RASTER_TARGET_WIDTH:
        scale = RASTER_TARGET_WIDTH / w_orig
        clip_img = clip_img.resize(
            (RASTER_TARGET_WIDTH, int(h_orig * scale)), Image.LANCZOS
        )
        dpi = SCREENSHOT_DPI * scale
    else:
        dpi = SCREENSHOT_DPI

    w_px, h_px = clip_img.size
    # スライドサイズ: 元の画像を SCREENSHOT_DPI で pt に換算
    w_pt = w_orig / SCREENSHOT_DPI * 72
    h_pt = h_orig / SCREENSHOT_DPI * 72

    print(f"処理開始: クリップボード画像 ({w_orig}x{h_orig}px)")
    if verbose:
        print(f"出力先:   {out_dir}")
        print(f"作業解像度: {w_px}x{h_px}px  仮想DPI={dpi:.1f}")
        print(f"スライドサイズ: {w_pt:.0f}x{h_pt:.0f}pt")
        print(f"OCR言語:  {OCR_LANGS}")

    det, rec = load_models(verbose=verbose)

    try:
        from pptx import Presentation  # noqa: F401
    except ImportError:
        print("[ERROR] python-pptx が必要です: pip install python-pptx")
        sys.exit(1)
    prs = _init_pptx(w_pt, h_pt)

    page_num = 0
    print(f"\n画像を処理中...")

    if verbose:
        print("  OCR実行中...")
    rows = extract_text(clip_img, dpi, page_num, det, rec, verbose=verbose)
    print(f"  テキスト: {len(rows)}行")

    masked = mask_text(clip_img, rows)
    del clip_img

    imgs = extract_blocks(masked, dpi, page_num, verbose=verbose)
    print(f"  画像ブロック: {len(imgs)}件")
    del masked

    _add_pptx_slide(prs, rows, imgs)

    if csv_path is not None:
        csv_imgs_list = [{k: v for k, v in m.items() if k != "img_bytes"} for m in imgs]
        write_meta_csv(rows, csv_imgs_list, csv_path)

    del rows, imgs
    gc.collect()

    prs.save(str(pptx_path))
    print(f"\nPowerPoint出力完了: {pptx_path}  (1スライド)")
    print("完了")


# ─────────────────────────────────────────────
# エントリーポイント
# ─────────────────────────────────────────────

def main() -> None:
    parser = argparse.ArgumentParser(
        description="PDF を PowerPoint に変換する (surya OCR使用)",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        "pdf", nargs="?", default=None,
        help="変換するPDFファイルのパス (省略時はクリップボードから画像またはパスを読み取り)",
    )
    parser.add_argument(
        "--csv", action="store_true", default=False,
        help="メタデータCSVを出力する (デフォルト: 出力しない)",
    )
    parser.add_argument(
        "-v", "--verbose", action="store_true", default=False,
        help="進捗を詳細表示する",
    )
    parser.add_argument(
        "--out-dir", default=None, metavar="DIR",
        help="出力先ディレクトリ (デフォルト: PDFと同じ場所)",
    )
    args = parser.parse_args()

    if args.pdf:
        # ── PDF パスを引数で受け取った場合 ──
        pdf_path = Path(args.pdf)
        out_dir = Path(args.out_dir) if args.out_dir else pdf_path.parent
        pptx_path = out_dir / (pdf_path.stem + ".pptx")
        csv_path = None
        if args.csv:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            csv_path = out_dir / f"pdf2pptx_meta_{ts}.csv"
        process_pdf(
            pdf_path, out_dir,
            pptx_path=pptx_path, csv_path=csv_path, verbose=args.verbose,
        )
    else:
        # ── 引数なし: クリップボードを確認 ──
        # 一時ファイルに出力して PowerPoint で直接開く（名前を付けて保存はユーザに委ねる）
        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        tmp.close()
        tmp_path = Path(tmp.name)

        csv_path = None
        if args.csv:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            csv_path = tmp_path.parent / f"pdf2pptx_meta_{ts}.csv"

        # 1) まず画像 (スクリーンショット等) を試みる
        clip_img = _read_clipboard_image()
        if clip_img is not None:
            print(f"クリップボード画像を取得: {clip_img.size[0]}x{clip_img.size[1]}px")
            process_image(
                clip_img, tmp_path.parent,
                pptx_path=tmp_path, csv_path=csv_path, verbose=args.verbose,
            )
        else:
            # 2) 画像がなければファイルパス (テキスト or HDROP) を試みる
            clip = _read_clipboard()
            if not clip:
                tmp_path.unlink(missing_ok=True)
                print("[ERROR] クリップボードが空です。PDFパスを引数で指定するか、"
                      "PDFパスまたはスクリーンショットをクリップボードにコピーしてください。")
                sys.exit(1)
            pdf_path = Path(clip)
            print(f"クリップボードから取得: {pdf_path}")
            process_pdf(
                pdf_path, tmp_path.parent,
                pptx_path=tmp_path, csv_path=csv_path, verbose=args.verbose,
            )

        print(f"PowerPoint で開いています...")
        os.startfile(str(tmp_path))


if __name__ == "__main__":
    main()
